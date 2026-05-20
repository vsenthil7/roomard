"""
Build the Roomard pitch deck for the Baidu "Build with MeDo" hackathon (AT-Hack0019).

12 slides, 16:9 widescreen, Roomard brand (deep teal #0a4a3f / pale-mint #e6f4f1 /
dark #073529). Editorial-style typography matching the SPA front-end (Inter).

Ported from the ATRIO Boardroom (AT-Hack0021) deck builder; same layout helpers,
re-themed to Roomard and re-written with Roomard's real content + metrics from:
  - README.md
  - docs/TRACEABILITY.md
  - docs/COVERAGE_BASELINE.md
  - docs/AT-Hack0019_Claude_Roomard_UseCaseCatalogue_*.md

Output (relative to roomard/ repo root):
  - submission_media/roomard-pitch-deck-{stamp}.pptx
  - submission_media/_backup/{same file}
  (then run scripts/pptx_to_pdf.ps1 to produce the .pdf deliverable)

This script is idempotent. Re-run to regenerate after content edits.
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ---------- Brand (Roomard) ----------
INK = RGBColor(0x03, 0x1A, 0x15)          # roomard-900 near-black teal background
PAPER = RGBColor(0xFF, 0xFF, 0xFF)        # text on dark / panel bg
TEXT_PRIMARY = RGBColor(0xE6, 0xF4, 0xF1) # roomard-50 body text on dark
TEXT_SECONDARY = RGBColor(0xBF, 0xE1, 0xD9)  # roomard-100 muted
TEAL = RGBColor(0x10, 0xB9, 0x81)         # accent (header strip / signal)
BLUE = RGBColor(0x3B, 0x82, 0xF6)         # scene-step accent
GREEN = RGBColor(0x10, 0xB9, 0x81)        # success
RED = RGBColor(0xDC, 0x26, 0x26)          # blocker / danger
PANEL_BG = RGBColor(0x07, 0x35, 0x29)     # roomard-700 panel against INK


# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HACK = "AT-HACK0019 \u00b7 BAIDU BUILD WITH MEDO 2026"
REPO = "github.com/vsenthil7/roomard"
TOTAL = 12


def fill_slide(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             italic=False, color=PAPER, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_top_strip(slide, color=TEAL):
    add_rect(slide, Inches(0.5), Inches(0.35), Inches(0.25), Inches(0.06), color)
    add_text(slide, Inches(0.85), Inches(0.25), Inches(11), Inches(0.3),
             HACK, size=10, color=TEXT_SECONDARY, font="Inter")


def add_footer(slide, page_num=None):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             f"Roomard \u00b7 {REPO} \u00b7 AI guest-memory engine",
             size=9, color=TEXT_SECONDARY, font="Inter")
    if page_num is not None:
        add_text(slide, Inches(11), Inches(7.05), Inches(2), Inches(0.3),
                 f"{page_num} / {TOTAL}", size=9, color=TEXT_SECONDARY,
                 font="Inter", align=PP_ALIGN.RIGHT)


# ---------- Slide builders ----------


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(1.5),
             "Roomard", size=92, bold=True, color=PAPER, font="Inter")
    add_rect(s, Inches(0.55), Inches(3.55), Inches(1.5), Inches(0.06), TEAL)
    add_text(s, Inches(0.5), Inches(3.85), Inches(12), Inches(1.0),
             "The AI guest-memory engine.",
             size=44, italic=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(4.85), Inches(12), Inches(1.0),
             "For boutique hotels. Capture every preference. Brief every shift. Prepare every room.",
             size=20, color=TEXT_PRIMARY, font="Inter")
    add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
             f"Baidu Build with MeDo 2026 \u00b7 AT-Hack0019 \u00b7 {REPO}",
             size=12, color=TEXT_SECONDARY, font="Inter")
    return s


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The problem", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Boutique hotels run on memory that walks out the door.",
             size=36, bold=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(3.5),
             "A returning guest's pillow preference, the allergy from last stay, the complaint that\n"
             "was never closed - all of it lives in one supervisor's head, a paper card in a drawer,\n"
             "or a concierge inbox nobody else can see.\n\n"
             "  \u2192 Staff turnover erases the guest relationship overnight.\n"
             "  \u2192 Chains have $200k CRMs; independents have sticky notes.\n"
             "  \u2192 The data exists - handwritten cards, emails, reviews, PMS records - but it is\n"
             "      never unified, never surfaced at the front desk, never auditable.",
             size=16, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 2)
    return s


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The solution", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "One memory engine. Three moments that matter.",
             size=34, bold=True, color=PAPER, font="Inter")
    cols = [
        ("CAPTURE", "Handwritten check-in cards, concierge emails, PMS records and reviews\nflow into one guest profile. PaddleOCR-VL + ERNIE 4.5 turn paper into structured, source-attributed preferences."),
        ("BRIEF", "Every morning, an AI arrival brief prioritises today's guests - VIP,\nrepeat, prior-issue - with one-tap evidence drill-down. Degrades gracefully when AI is unavailable."),
        ("PREPARE", "Mid-conversation lookup, D-1 housekeeping prep cards, and a\ncomplaint-trajectory flag that catches a worsening guest before the next stay. Every action audited."),
    ]
    col_w = Inches(4.1)
    gap = Inches(0.1)
    for i, (header, body) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        add_rect(s, x, Inches(2.8), col_w, Inches(0.06), TEAL)
        add_text(s, x, Inches(2.95), col_w, Inches(0.4),
                 header, size=14, bold=True, color=TEAL, font="Inter")
        add_text(s, x, Inches(3.45), col_w, Inches(3),
                 body, size=13, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 3)
    return s


def slide_use_cases(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The MVP wedge", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Eight use cases. Seven built. One honest stub.",
             size=30, bold=True, color=PAPER, font="Inter")
    rows = [
        ("UC-07", "Daily arrival brief", "90%", GREEN),
        ("UC-01", "Handwritten card capture (PaddleOCR-VL)", "85%", GREEN),
        ("UC-08", "Mid-conversation guest lookup", "80%", GREEN),
        ("UC-09", "Housekeeping prep card (D-1)", "80%", GREEN),
        ("UC-24a", "Mews inbound PMS sync", "80%", GREEN),
        ("UC-23", "Confidence / exception queue", "75%", GREEN),
        ("UC-11", "Complaint trajectory (ERNIE X1) - bonus", "75%", TEAL),
        ("UC-25", "Review polling (pipeline real, APIs stubbed)", "70%", BLUE),
        ("UC-29", "SSO (SAML/OIDC) - honest 501", "5%", RED),
    ]
    y0 = Inches(2.6)
    for i, (uc, name, pct, color) in enumerate(rows):
        y = y0 + Inches(0.46) * i
        add_text(s, Inches(0.5), y, Inches(1.2), Inches(0.4),
                 uc, size=13, bold=True, color=TEAL, font="Inter")
        add_text(s, Inches(1.9), y, Inches(8.5), Inches(0.4),
                 name, size=13, color=TEXT_PRIMARY, font="Inter")
        add_text(s, Inches(10.6), y, Inches(2.0), Inches(0.4),
                 pct, size=13, bold=True, color=color, font="Inter", align=PP_ALIGN.RIGHT)
    add_footer(s, 4)
    return s


def slide_ai(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The AI stack", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Sponsor models, routed through one gateway.",
             size=30, bold=True, color=PAPER, font="Inter")
    items = [
        ("PaddleOCR-VL", "Layout-aware OCR", "Handwritten check-in cards \u2192 text + structure"),
        ("ERNIE 4.5", "Entity extraction", "Card/email/review text \u2192 structured preferences"),
        ("ERNIE X1", "Reasoning", "Complaint-trajectory flag (3-issue rule, UC-11)"),
        ("Qianfan MaaS", "Inference gateway", "Single path to models; per-tenant caps; every call logged"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(1.2)
    for i, (name, role, detail) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + (col_w + Inches(0.3)) * col
        y = Inches(2.8) + (row_h + Inches(0.3)) * row
        add_rect(s, x, y, col_w, row_h, PANEL_BG, line_color=TEAL)
        add_text(s, x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.5), Inches(0.4),
                 name, size=18, bold=True, color=PAPER, font="Inter")
        add_text(s, x + Inches(0.25), y + Inches(0.52), col_w - Inches(0.5), Inches(0.3),
                 role, size=12, color=TEAL, font="Inter")
        add_text(s, x + Inches(0.25), y + Inches(0.82), col_w - Inches(0.5), Inches(0.3),
                 detail, size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9),
             "AI degrades gracefully everywhere: when a model is unavailable the brief still renders, the prep\n"
             "card falls back to a plain card, and the trajectory flag falls back to the rule-based threshold.",
             size=12, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 5)
    return s


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Live demo \u00b7 4 stages", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "What you will see in the walkthrough.",
             size=30, bold=True, color=PAPER, font="Inter")
    stages = [
        ("1", "DAILY ARRIVAL BRIEF (UC-07)",
         "Front-desk manager signs in and lands on the morning brief: prioritised arrivals, AI-written notes, one-tap evidence."),
        ("2", "GUEST LOOKUP + TRAJECTORY (UC-08 / UC-11)",
         "Search a guest mid-conversation: 3-bullet preferences, last-stay summary, 'say this' line, complaint-trajectory flag."),
        ("3", "HANDWRITTEN CARD CAPTURE (UC-01)",
         "Photograph a paper check-in card; PaddleOCR-VL + ERNIE 4.5 extract fields; low confidence routes to the queue."),
        ("4", "EXCEPTION QUEUE + PREP CARDS (UC-23 / UC-09)",
         "Clear the confidence queue; review tomorrow's housekeeping prep cards; two-tap completion; every action audited."),
    ]
    for i, (num, title, body) in enumerate(stages):
        y = Inches(2.7) + Inches(1.05) * i
        add_text(s, Inches(0.5), y, Inches(0.8), Inches(1),
                 num, size=44, bold=True, color=TEAL, font="Inter")
        add_text(s, Inches(1.4), y + Inches(0.1), Inches(11.4), Inches(0.4),
                 title, size=15, bold=True, color=PAPER, font="Inter")
        add_text(s, Inches(1.4), y + Inches(0.5), Inches(11.4), Inches(0.6),
                 body, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 6)
    return s


def slide_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Proof, not promises", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Engineering rigour you can re-run.",
             size=30, bold=True, color=PAPER, font="Inter")
    metrics = [
        ("356", "unit tests pass"),
        ("12", "DB integration tests (live PG)"),
        ("87.5 %", "mean module coverage"),
        ("39", "findings fixed \u00b7 0 open"),
        ("3", "production bugs caught by tests"),
        ("0 / 0", "lint errors / warnings"),
        ("78", "traceable commits on main"),
        ("RLS", "tenant isolation proven live"),
    ]
    box_w = Inches(3.0)
    box_h = Inches(1.4)
    for i, (big, small) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + (box_w + Inches(0.1)) * col
        y = Inches(2.9) + (box_h + Inches(0.2)) * row
        add_rect(s, x, y, box_w, box_h, PANEL_BG, line_color=GREEN)
        add_text(s, x, y + Inches(0.2), box_w, Inches(0.6),
                 big, size=28, bold=True, color=PAPER, font="Inter", align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.85), box_w, Inches(0.45),
                 small, size=11, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
             "Reproducible: pnpm -r test  \u00b7  pnpm -r run test:coverage  \u00b7  DB integration gated on DATABASE_URL",
             size=11, color=TEXT_SECONDARY, font="Consolas", italic=True)
    add_footer(s, 7)
    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Architecture", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "A typed monorepo. Ten services. One edge.",
             size=30, bold=True, color=PAPER, font="Inter")
    tiers = [
        ("FRONTEND", "React 18 \u00b7 TanStack Router \u00b7 Vite \u00b7 Tailwind (PWA)", "Offline capture queue \u00b7 served on :8180", BLUE),
        ("EDGE + SERVICES", "API gateway + 10 Fastify services (auth, guest, capture, brief, exception, audit, ingest, tenant, ai-gateway)", "JWT + RBAC at the edge \u00b7 gateway on :3100", TEAL),
        ("DATA + AI", "Postgres 16 (RLS-enforced) \u00b7 Redis \u00b7 MinIO \u00b7 Qianfan MaaS", "Append-only audit \u00b7 per-tenant row-level security", GREEN),
    ]
    for i, (label, stack, note, color) in enumerate(tiers):
        y = Inches(2.8) + Inches(1.2) * i
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), PANEL_BG, line_color=color)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(3.0), Inches(0.4),
                 label, size=15, bold=True, color=color, font="Inter")
        add_text(s, Inches(4.0), y + Inches(0.1), Inches(8.6), Inches(0.4),
                 stack, size=12, bold=True, color=PAPER, font="Inter")
        add_text(s, Inches(4.0), y + Inches(0.55), Inches(8.6), Inches(0.4),
                 note, size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 8)
    return s


def slide_trust(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Why this is trustworthy", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Multi-tenant and audit-grade from line one.",
             size=30, bold=True, color=PAPER, font="Inter")
    points = [
        ("Row-level security, proven live", "RLS is forced on every tenant table; a dedicated non-superuser app role makes isolation real - demonstrated against the live database, not just claimed."),
        ("Append-only audit chain", "Every event is hash-chained; the verify endpoint re-derives the chain in SQL so tampering is detectable."),
        ("Tests find real bugs", "Three production bugs (a tenant-context failure, an audit schema drift, a dead error-forwarder) were caught by chasing real coverage - not padding."),
        ("Graceful degradation", "Every AI path has a non-AI fallback; the product stays useful when a model is down."),
    ]
    for i, (head, body) in enumerate(points):
        y = Inches(2.7) + Inches(1.05) * i
        add_rect(s, Inches(0.5), y + Inches(0.05), Inches(0.06), Inches(0.85), TEAL)
        add_text(s, Inches(0.75), y, Inches(12), Inches(0.4),
                 head, size=15, bold=True, color=PAPER, font="Inter")
        add_text(s, Inches(0.75), y + Inches(0.42), Inches(12), Inches(0.6),
                 body, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 9)
    return s


def slide_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "What ships next", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "From wedge to platform.",
             size=30, bold=True, color=PAPER, font="Inter")
    columns = [
        ("Close the wedge", [
            "Real review-API adapters (TripAdvisor / Booking / Google)",
            "Outbound PMS sync (UC-24b) + a second PMS connector",
            "SSO via a real IdP test tenant (replace the 501)",
            "Push-notification delivery for briefs + GM flags",
            "Benchmark corpora to verify the OCR / ranking AC bars",
        ]),
        ("Platform", [
            "GDPR subject-access + erasure (EU legal gate)",
            "Native mobile (PWA \u2192 Capacitor) for housekeeping",
            "Multi-language UI (EN / IT / ES / ZH)",
            "Long-horizon guest memory across properties",
            "Lock the live stack on the non-superuser role end-to-end",
        ]),
    ]
    for i, (header, items) in enumerate(columns):
        x = Inches(0.5) + Inches(6.4) * i
        add_rect(s, x, Inches(2.7), Inches(6.0), Inches(0.05), TEAL)
        add_text(s, x, Inches(2.85), Inches(6.0), Inches(0.4),
                 header, size=15, bold=True, color=PAPER, font="Inter")
        for j, item in enumerate(items):
            add_text(s, x, Inches(3.35) + Inches(0.5) * j, Inches(6.0), Inches(0.4),
                     "\u00b7  " + item, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 10)
    return s


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Team + ask", size=12, color=TEAL, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Built for the Baidu Build with MeDo hackathon.",
             size=28, bold=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(2.8), Inches(12), Inches(0.4),
             "Founder + lead engineer \u00b7 product, architecture, full-stack build, demo, verification",
             size=14, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(3.25), Inches(12), Inches(0.4),
             "Generated under the AT-Hack0019 sprint with Claude as paired engineer.",
             size=12, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_text(s, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
             "What we are asking for",
             size=16, bold=True, color=TEAL, font="Inter")
    asks = [
        "Qianfan production-tier access for the live multilingual rollout.",
        "A boutique-hotel design partner (10-50 keys) for a 4-week pilot.",
        "A Mews / Cloudbeds partner connection to validate the PMS sync at scale.",
        "Time to demo the capture \u2192 brief \u2192 prepare loop with a real property's data.",
    ]
    for i, ask in enumerate(asks):
        add_text(s, Inches(0.5), Inches(5.0) + Inches(0.4) * i, Inches(12), Inches(0.4),
                 "\u00b7  " + ask, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 11)
    return s


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(1.4),
             "Roomard", size=72, bold=True, color=PAPER, font="Inter")
    add_rect(s, Inches(0.55), Inches(3.5), Inches(1.5), Inches(0.06), TEAL)
    add_text(s, Inches(0.5), Inches(3.75), Inches(12), Inches(1.0),
             "Every preference remembered.",
             size=28, italic=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(4.4), Inches(12), Inches(1.0),
             "Capture. Brief. Prepare. Multi-tenant, RLS-enforced, audit-grade by default.",
             size=18, color=TEXT_PRIMARY, font="Inter")
    add_text(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
             REPO, size=14, color=TEAL, font="Consolas")
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
             "356 tests \u00b7 12 DB integration \u00b7 87.5% mean coverage \u00b7 39 findings fixed \u00b7 0 open",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "Thank you.", size=12, color=TEXT_SECONDARY, font="Inter")
    return s


# ---------- Main ----------


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_use_cases,
        slide_ai,
        slide_demo,
        slide_proof,
        slide_architecture,
        slide_trust,
        slide_next,
        slide_team,
        slide_closing,
    ]
    for b in builders:
        b(prs)

    out_dir = Path(__file__).resolve().parent.parent / "submission_media"
    out_dir.mkdir(parents=True, exist_ok=True)
    backup_dir = out_dir / "_backup"
    backup_dir.mkdir(exist_ok=True)

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_pptx = out_dir / f"roomard-pitch-deck-{stamp}.pptx"
    prs.save(out_pptx)
    import shutil
    shutil.copy2(out_pptx, backup_dir / out_pptx.name)

    print(f"deck: {out_pptx} ({out_pptx.stat().st_size / 1024:.1f} KB)")
    return out_pptx


if __name__ == "__main__":
    build_deck()
